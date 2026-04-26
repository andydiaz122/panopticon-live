#!/usr/bin/env python3
"""
build_card_03_closing.py
========================

Generates Card 3 (closing thesis) for the PANOPTICON LIVE demo video, in
TWO formats from a single spec:

    1. card_03_closing.pptx — editable PowerPoint slide (open in PowerPoint
       to tweak; export-to-PNG via PowerPoint preserves spec exactly)
    2. card_03_closing.png  — final PNG ready to drop into CapCut directly
       (no PowerPoint required for the actual demo assembly path)

Spec source-of-truth:
    demo-presentation/scripts/title_card_specs.md § "✅ ACTIVE — Card 3"

Why this script over Keynote UI:
    - Deterministic output (no hex-color picker drift, no kerning slider creep)
    - Version-controlled (commit-able, reproducible across machines)
    - Single-source-of-truth (the spec lives in this script, not split between
      a Keynote file + a markdown spec that can drift)
    - Fast iteration (re-run after spec change, ~1 second; vs Keynote 5-10 min
      to rebuild)
    - Honors team-lead "Anthropic Minimalism" + Andrew's "use Anthropic skills"
      directive (uses python-pptx from the official anthropics/skills repo
      installed at ~/.claude/skills/pptx/)

Run:
    /Users/andrew/Documents/Coding/Built_with_Opus_4-7_Hackathon/.venv/bin/python \
        scripts/build_card_03_closing.py

Output:
    ~/Documents/Panopticon_TitleCards/card_03_closing.pptx
    ~/Documents/Panopticon_TitleCards/card_03_closing.png
"""

from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

# ============ Spec ============

# Slide dimensions (Keynote/PowerPoint standard 16:9 at 1920×1080 px)
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080

# Pt → Px conversion at 96 DPI (standard screen render, matches Keynote's
# native 1080p export). 1pt = 96/72 = 1.333 px.
PT_TO_PX = 96 / 72

# Colors (all hex from title_card_specs.md spec)
BG_HEX        = "#000000"     # pure void background
INK_HEX       = "#F8FAFC"     # near-white headline (NOT pure white)
ACCENT_HEX    = "#00E5FF"     # cyan accent rule (NOT Apple's brand blue)
WHISPER_HEX   = "#5A6678"     # cool muted chrome (built-with + URL)

# Headline (Fraunces serif, 56pt, near-white, -0.5% letter-spacing, centered)
HEADLINE_TEXT = "capture the signal nobody else is reading."
HEADLINE_PT   = 56
HEADLINE_TRACKING = -0.005   # -0.5% as a unitless ratio

# Cyan accent rule (1pt × 280pt, 70% opacity, centered, 56pt below headline)
RULE_WIDTH_PT  = 280
RULE_HEIGHT_PT = 1
RULE_OPACITY   = 0.70
RULE_GAP_BELOW_HEADLINE_PT = 56

# Built-with line (JetBrains Mono, 13pt, whisper, ALL CAPS, +24% spacing)
BUILTWITH_TEXT = "built with claude opus 4.7"
BUILTWITH_PT   = 13
BUILTWITH_TRACKING = 0.24
BUILTWITH_GAP_BELOW_RULE_PT = 80

# GitHub URL (JetBrains Mono, 13pt, whisper, lowercase, +6% spacing)
GITHUB_TEXT = "github.com/andydiaz122/panopticon-live"
GITHUB_PT   = 13
GITHUB_TRACKING = 0.06
GITHUB_GAP_BELOW_BUILTWITH_PT = 14

# Output paths
OUT_DIR = Path.home() / "Documents" / "Panopticon_TitleCards"
OUT_PPTX = OUT_DIR / "card_03_closing.pptx"
OUT_PNG  = OUT_DIR / "card_03_closing.png"

# Font paths — use STATIC Regular (400) weights, NOT the variable fonts.
# Pillow's variable-font support is iffy and defaults to a heavier weight
# than Regular, producing chunky-looking serifs. The static Regular files
# live in the Fraunces zip Andrew downloaded; we point at them directly so
# the script is reproducible regardless of system font install state.
FRAUNCES_TTF = Path.home() / "Downloads" / "Fraunces" / "static" / "Fraunces_72pt-Regular.ttf"
# JetBrains Mono variable font does default to Regular cleanly so we keep it.
JETBRAINS_TTF = Path.home() / "Library" / "Fonts" / "JetBrainsMono-VariableFont_wght.ttf"
# Fallbacks if Andrew's Downloads layout differs:
FRAUNCES_FALLBACKS = [
    Path.home() / "Downloads" / "Fraunces" / "static" / "Fraunces_72pt_Soft-Regular.ttf",
    Path.home() / "Library" / "Fonts" / "Fraunces-Regular.ttf",
    Path.home() / "Library" / "Fonts" / "Fraunces-VariableFont_SOFT,WONK,opsz,wght.ttf",
]


# ============ Helpers ============

def _hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
    s = hex_str.lstrip("#")
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _hex_to_rgba(hex_str: str, alpha: float) -> tuple[int, int, int, int]:
    r, g, b = _hex_to_rgb(hex_str)
    return (r, g, b, int(alpha * 255))


def pt_to_px(pt: float) -> int:
    """Convert pt (PowerPoint/Keynote unit) → px (PNG render unit)."""
    return round(pt * PT_TO_PX)


# ============ PPTX builder ============

def build_pptx() -> None:
    """Build the editable .pptx via python-pptx."""
    prs = Presentation()
    # Set slide size to 1920×1080 px (16:9). python-pptx uses EMU
    # (English Metric Units): 1 pt = 12700 EMU, 1 inch = 914400 EMU.
    # 1920px @ 96 DPI = 20" = 18288000 EMU; 1080px @ 96 DPI = 11.25" = 10287000 EMU.
    prs.slide_width = Emu(int(SLIDE_W_PX / 96 * 914400))
    prs.slide_height = Emu(int(SLIDE_H_PX / 96 * 914400))

    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Background (full-bleed black rectangle)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(BG_HEX))
    bg_shape.line.fill.background()  # no border

    # Compute vertical center group: estimate group height via text heights
    # Using rough vertical-rhythm assumptions matching the px renderer below.
    headline_h_pt = HEADLINE_PT * 1.2  # line-height ~1.2
    builtwith_h_pt = BUILTWITH_PT * 1.2
    github_h_pt = GITHUB_PT * 1.2
    group_height_pt = (
        headline_h_pt
        + RULE_GAP_BELOW_HEADLINE_PT
        + RULE_HEIGHT_PT
        + BUILTWITH_GAP_BELOW_RULE_PT
        + builtwith_h_pt
        + GITHUB_GAP_BELOW_BUILTWITH_PT
        + github_h_pt
    )
    # Convert pt → EMU and center vertically
    group_height_emu = int(group_height_pt * 12700)
    group_top_emu = (prs.slide_height - group_height_emu) // 2
    cursor_top_emu = group_top_emu

    # Headline text box (centered horizontally, full slide width)
    headline_h_emu = int(headline_h_pt * 12700)
    headline_box = slide.shapes.add_textbox(
        Emu(0), Emu(cursor_top_emu), prs.slide_width, Emu(headline_h_emu)
    )
    tf = headline_box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = 2  # center (PP_ALIGN.CENTER = 2)
    run = p.add_run()
    run.text = HEADLINE_TEXT
    run.font.name = "Fraunces"
    run.font.size = Pt(HEADLINE_PT)
    run.font.color.rgb = RGBColor(*_hex_to_rgb(INK_HEX))
    # Kerning / letter-spacing: python-pptx exposes spc on a:rPr. We set via XML.
    rPr = run.font._rPr
    spc_val = int(HEADLINE_TRACKING * HEADLINE_PT * 100)  # in centi-points, signed
    rPr.set("spc", str(spc_val))
    cursor_top_emu += headline_h_emu + int(RULE_GAP_BELOW_HEADLINE_PT * 12700)

    # Cyan accent rule (a thin filled rectangle, centered horizontally)
    rule_w_emu = int(RULE_WIDTH_PT * 12700)
    rule_h_emu = int(RULE_HEIGHT_PT * 12700)
    rule_left_emu = (prs.slide_width - rule_w_emu) // 2
    rule_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(rule_left_emu),
        Emu(cursor_top_emu),
        Emu(rule_w_emu),
        Emu(rule_h_emu),
    )
    rule_shape.fill.solid()
    rule_shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(ACCENT_HEX))
    # Apply 70% opacity via XML (python-pptx doesn't expose opacity on Fill)
    fill_xml = rule_shape.fill._xPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
    )
    srgbClr = fill_xml.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
    )
    from lxml import etree
    alpha = etree.SubElement(
        srgbClr, "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha"
    )
    alpha.set("val", str(int(RULE_OPACITY * 100000)))  # OpenXML uses thousandths of a percent
    rule_shape.line.fill.background()
    cursor_top_emu += rule_h_emu + int(BUILTWITH_GAP_BELOW_RULE_PT * 12700)

    # Built-with line (JetBrains Mono, ALL CAPS via uppercase string)
    builtwith_h_emu = int(builtwith_h_pt * 12700)
    builtwith_box = slide.shapes.add_textbox(
        Emu(0), Emu(cursor_top_emu), prs.slide_width, Emu(builtwith_h_emu)
    )
    tf = builtwith_box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = 2
    run = p.add_run()
    run.text = BUILTWITH_TEXT.upper()
    run.font.name = "JetBrains Mono"
    run.font.size = Pt(BUILTWITH_PT)
    run.font.color.rgb = RGBColor(*_hex_to_rgb(WHISPER_HEX))
    rPr = run.font._rPr
    rPr.set("spc", str(int(BUILTWITH_TRACKING * BUILTWITH_PT * 100)))
    cursor_top_emu += builtwith_h_emu + int(GITHUB_GAP_BELOW_BUILTWITH_PT * 12700)

    # GitHub URL line (JetBrains Mono, lowercase)
    github_h_emu = int(github_h_pt * 12700)
    github_box = slide.shapes.add_textbox(
        Emu(0), Emu(cursor_top_emu), prs.slide_width, Emu(github_h_emu)
    )
    tf = github_box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = 2
    run = p.add_run()
    run.text = GITHUB_TEXT
    run.font.name = "JetBrains Mono"
    run.font.size = Pt(GITHUB_PT)
    run.font.color.rgb = RGBColor(*_hex_to_rgb(WHISPER_HEX))
    rPr = run.font._rPr
    rPr.set("spc", str(int(GITHUB_TRACKING * GITHUB_PT * 100)))

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"✓ PPTX written: {OUT_PPTX}")


# ============ PNG renderer (Pillow) ============

def _measure_text_with_tracking(
    text: str, font: ImageFont.FreeTypeFont, tracking_ratio: float
) -> int:
    """Measure rendered width of text with per-character tracking applied."""
    if not text:
        return 0
    base_size = font.size
    extra_per_char = round(tracking_ratio * base_size)
    total = 0
    for ch in text:
        bbox = font.getbbox(ch)
        ch_width = bbox[2] - bbox[0]
        total += ch_width + extra_per_char
    total -= extra_per_char  # no trailing extra after last char
    return total


def _draw_text_with_tracking(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    font: ImageFont.FreeTypeFont,
    fill: tuple[int, int, int],
    tracking_ratio: float,
) -> None:
    """Draw text character-by-character with custom letter-spacing."""
    x, y = xy
    extra = round(tracking_ratio * font.size)
    for ch in text:
        draw.text((x, y), ch, font=font, fill=fill)
        bbox = font.getbbox(ch)
        ch_width = bbox[2] - bbox[0]
        x += ch_width + extra


def _resolve_fraunces_font() -> Path:
    """Find a usable Fraunces TTF — prefer static Regular, fall back through known paths."""
    candidates = [FRAUNCES_TTF] + FRAUNCES_FALLBACKS
    for p in candidates:
        if p.exists():
            return p
    raise FileNotFoundError(f"No Fraunces font found in: {[str(c) for c in candidates]}")


def _measure_rendered_bbox(
    text: str, font: ImageFont.FreeTypeFont, tracking_ratio: float
) -> tuple[int, int]:
    """
    Measure rendered (width, height) of text with tracking. Height comes from
    the font's metric ascent+descent so vertical centering accounts for full
    rendered extent (including descenders), not just cap-height.
    """
    width = _measure_text_with_tracking(text, font, tracking_ratio)
    # Pillow's getmetrics returns (ascent, descent); height = ascent + descent.
    ascent, descent = font.getmetrics()
    height = ascent + descent
    return width, height


def build_png() -> None:
    """Build the final PNG via Pillow with exact pixel positioning."""
    fraunces_path = _resolve_fraunces_font()
    print(f"  Using Fraunces from: {fraunces_path}")
    if not JETBRAINS_TTF.exists():
        print(f"✗ JetBrains Mono font not found at {JETBRAINS_TTF}", file=sys.stderr)
        sys.exit(1)

    img = Image.new("RGBA", (SLIDE_W_PX, SLIDE_H_PX), _hex_to_rgba(BG_HEX, 1.0))
    draw = ImageDraw.Draw(img, "RGBA")

    # Load fonts at pixel-equivalent sizes
    headline_font = ImageFont.truetype(str(fraunces_path), pt_to_px(HEADLINE_PT))
    builtwith_font = ImageFont.truetype(str(JETBRAINS_TTF), pt_to_px(BUILTWITH_PT))
    github_font = ImageFont.truetype(str(JETBRAINS_TTF), pt_to_px(GITHUB_PT))

    # Measure rendered (width, height) using REAL font metrics (ascent + descent)
    # — this fixes the upper-half-skew bug from the v1 render where headline_h
    # was estimated as `pt_to_px(pt * 1.0)` which under-counted the rendered
    # vertical extent by ~30% (no descender accounting).
    headline_w, headline_h = _measure_rendered_bbox(
        HEADLINE_TEXT, headline_font, HEADLINE_TRACKING
    )
    builtwith_w, builtwith_h = _measure_rendered_bbox(
        BUILTWITH_TEXT.upper(), builtwith_font, BUILTWITH_TRACKING
    )
    github_w, github_h = _measure_rendered_bbox(
        GITHUB_TEXT, github_font, GITHUB_TRACKING
    )

    rule_h = max(1, pt_to_px(RULE_HEIGHT_PT))
    rule_w = pt_to_px(RULE_WIDTH_PT)
    gap_rule = pt_to_px(RULE_GAP_BELOW_HEADLINE_PT)
    gap_builtwith = pt_to_px(BUILTWITH_GAP_BELOW_RULE_PT)
    gap_github = pt_to_px(GITHUB_GAP_BELOW_BUILTWITH_PT)

    group_total_h = (
        headline_h + gap_rule + rule_h + gap_builtwith + builtwith_h + gap_github + github_h
    )
    cursor_y = (SLIDE_H_PX - group_total_h) // 2

    # Draw headline (Fraunces serif, near-white, centered, with -0.5% tracking)
    headline_x = (SLIDE_W_PX - headline_w) // 2
    _draw_text_with_tracking(
        draw, (headline_x, cursor_y), HEADLINE_TEXT, headline_font,
        _hex_to_rgb(INK_HEX), HEADLINE_TRACKING,
    )
    cursor_y += headline_h + gap_rule

    # Draw cyan accent rule (with 70% opacity)
    rule_x = (SLIDE_W_PX - rule_w) // 2
    draw.rectangle(
        [(rule_x, cursor_y), (rule_x + rule_w, cursor_y + rule_h)],
        fill=_hex_to_rgba(ACCENT_HEX, RULE_OPACITY),
    )
    cursor_y += rule_h + gap_builtwith

    # Draw built-with line (JetBrains Mono, ALL CAPS, +24% tracking)
    builtwith_x = (SLIDE_W_PX - builtwith_w) // 2
    _draw_text_with_tracking(
        draw, (builtwith_x, cursor_y), BUILTWITH_TEXT.upper(), builtwith_font,
        _hex_to_rgb(WHISPER_HEX), BUILTWITH_TRACKING,
    )
    cursor_y += builtwith_h + gap_github

    # Draw GitHub URL (JetBrains Mono, lowercase, +6% tracking)
    github_x = (SLIDE_W_PX - github_w) // 2
    _draw_text_with_tracking(
        draw, (github_x, cursor_y), GITHUB_TEXT, github_font,
        _hex_to_rgb(WHISPER_HEX), GITHUB_TRACKING,
    )

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    img.convert("RGB").save(str(OUT_PNG), "PNG", optimize=True)
    print(f"✓ PNG written:  {OUT_PNG}  ({OUT_PNG.stat().st_size:,} bytes, {SLIDE_W_PX}×{SLIDE_H_PX})")


# ============ Main ============

def main() -> int:
    print("Building Card 3 — closing thesis (PANOPTICON LIVE)")
    print(f"  Spec source: demo-presentation/scripts/title_card_specs.md")
    print(f"  Headline: {HEADLINE_TEXT!r}")
    print(f"  Output dir: {OUT_DIR}")
    print()
    build_pptx()
    build_png()
    print()
    print("✓ DONE — Card 3 generated in both formats.")
    print()
    print("Next steps:")
    print(f"  1. Inspect PNG: open {OUT_PNG}")
    print(f"  2. (Optional) Open PPTX in PowerPoint: open {OUT_PPTX}")
    print(f"  3. Drop card_03_closing.png into CapCut Media bin")
    return 0


if __name__ == "__main__":
    sys.exit(main())
